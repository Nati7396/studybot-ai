import os
import logging
import re
import pdfplumber
import PyPDF2
from bot.config import UPLOAD_DIR, MAX_CHUNK_SIZE, CHUNK_OVERLAP

logger = logging.getLogger(__name__)


def extract_text_from_file(file_path: str) -> tuple[str, int]:
    """Route extraction by file type. Returns (text, page_count)."""
    ext = os.path.splitext(file_path.lower())[1]
    if ext == ".pdf":
        text, pages = extract_text_from_pdf(file_path)
        if not text or len(text) < 50:
            logger.info("Text PDF extraction failed — trying Gemini Vision OCR...")
            text, pages = extract_text_with_gemini_ocr(file_path)
        return text, pages
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    else:
        return "", 0


def extract_text_from_pdf(file_path: str) -> tuple[str, int]:
    """Extract text from PDF using pdfplumber then PyPDF2."""
    text_parts = []
    page_count = 0
    try:
        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                try:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
                except Exception as e:
                    logger.warning(f"pdfplumber failed on page {i}: {e}")
    except Exception as e:
        logger.warning(f"pdfplumber failed entirely: {e}, trying PyPDF2")

    if not text_parts:
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                page_count = len(reader.pages)
                for i, page in enumerate(reader.pages):
                    try:
                        t = page.extract_text()
                        if t:
                            text_parts.append(t)
                    except Exception as e:
                        logger.warning(f"PyPDF2 failed on page {i}: {e}")
        except Exception as e:
            logger.error(f"Both PDF extractors failed: {e}")

    full_text = clean_text("\n\n".join(text_parts))
    return full_text, page_count


def extract_text_with_gemini_ocr(file_path: str) -> tuple[str, int]:
    """Use Gemini Vision to OCR a scanned/image-based PDF."""
    try:
        import google.generativeai as genai
        from bot.config import GEMINI_API_KEYS
        if not GEMINI_API_KEYS:
            return "", 0
        genai.configure(api_key=GEMINI_API_KEYS[0])
        logger.info(f"Uploading {file_path} to Gemini for OCR...")
        uploaded = genai.upload_file(file_path, mime_type="application/pdf")
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content([
            "Extract ALL text from this document exactly as it appears. "
            "Include every word, heading, bullet point, number, formula, and table. "
            "Output ONLY the raw extracted text — no commentary, no markdown formatting.",
            uploaded
        ])
        try:
            genai.delete_file(uploaded.name)
        except Exception:
            pass
        text = clean_text(response.text or "")
        return text, 0
    except Exception as e:
        logger.error(f"Gemini OCR failed: {e}")
        return "", 0


def extract_text_from_pptx(file_path: str) -> tuple[str, int]:
    """Extract text from PowerPoint .pptx files."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        slide_count = len(prs.slides)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        full_text = clean_text("\n\n".join(text_parts))
        return full_text, slide_count
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return "", 0


def extract_text_from_docx(file_path: str) -> tuple[str, int]:
    """Extract text from Word .docx files."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        full_text = clean_text("\n\n".join(text_parts))
        return full_text, 1
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return "", 0


def extract_text_from_txt(file_path: str) -> tuple[str, int]:
    """Extract text from plain text files."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return clean_text(text), 1
    except Exception as e:
        logger.error(f"TXT extraction failed: {e}")
        return "", 0


def clean_text(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    return text.strip()


def split_into_chunks(text: str, chunk_size: int = MAX_CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks by sentence boundaries."""
    if not text:
        return []
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_len = 0
    for sentence in sentences:
        slen = len(sentence)
        if current_len + slen > chunk_size and current_chunk:
            chunks.append(" ".join(current_chunk))
            overlap_words = " ".join(current_chunk).split()[-overlap // 10:]
            current_chunk = [" ".join(overlap_words)]
            current_len = len(current_chunk[0])
        current_chunk.append(sentence)
        current_len += slen + 1
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    return [c for c in chunks if len(c.strip()) > 50]


def save_uploaded_file(user_id: int, file_bytes: bytes, original_name: str) -> str:
    user_dir = os.path.join(UPLOAD_DIR, str(user_id))
    os.makedirs(user_dir, exist_ok=True)
    safe_name = re.sub(r'[^\w\-_.]', '_', original_name)
    file_path = os.path.join(user_dir, safe_name)
    with open(file_path, "wb") as f:
        f.write(file_bytes)
    return file_path
